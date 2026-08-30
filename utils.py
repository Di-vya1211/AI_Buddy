import io
import os

import streamlit as st
from groq import APIConnectionError, APIStatusError, Groq, RateLimitError
from sklearn.feature_extraction.text import TfidfVectorizer

MODEL = "openai/gpt-oss-20b"

# ---------------------------------------------------------------------------
# Groq client
# ---------------------------------------------------------------------------

def get_client():
    """Return an authenticated Groq client, or None if no API key is set.

    The key is looked up first in Streamlit session state (sidebar input),
    then falls back to the GROQ_API_KEY environment variable loaded from .env.
    Returns None (not raises) so callers can give the user a friendly prompt.
    """
    api_key = st.session_state.get("api_key") or os.environ.get("GROQ_API_KEY")
    if not api_key:
        return None
    return Groq(api_key=api_key)


def ask(prompt, json_mode=False):
    """Send *prompt* to the Groq chat completion endpoint and return the text.

    Parameters
    ----------
    prompt : str
        The full prompt text to send as a single user message.
    json_mode : bool
        When True, sets ``response_format={"type": "json_object"}`` so the
        model is constrained to return valid JSON.  Use for all structured
        outputs (quiz, flashcards, concept map, Feynman feedback).

    Returns
    -------
    str
        The assistant's response, stripped of leading/trailing whitespace.

    Raises
    ------
    RuntimeError
        If the API key is missing, the request is rate-limited, or any other
        API-level error occurs — with a human-friendly message suitable for
        direct display via ``st.error()``.
    """
    client = get_client()
    if client is None:
        raise RuntimeError(
            "Groq API key is not set. Paste your key in the sidebar, or add it "
            "to a `.env` file as `GROQ_API_KEY=...` (see README)."
        )

    kwargs = {}
    if json_mode:
        kwargs["response_format"] = {"type": "json_object"}

    try:
        completion = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=4096,  # raised from 2048 — quiz + feedback can be long
            **kwargs,
        )
    except RateLimitError:
        raise RuntimeError(
            "Rate limit reached on your Groq key. Wait a moment and try again, "
            "or check your usage at console.groq.com."
        )
    except APIConnectionError:
        raise RuntimeError(
            "Could not reach the Groq API. Check your internet connection and try again."
        )
    except APIStatusError as exc:
        # Covers 4xx/5xx responses not handled above (e.g. invalid model name)
        raise RuntimeError(f"Groq API error ({exc.status_code}): {exc.message}")

    return completion.choices[0].message.content.strip()


def ask_stream(prompt):
    """Stream a non-JSON response from Groq; yields text chunks as they arrive.

    This is used for free-text tabs (Explain, Revision plan, Feynman plain
    text) so tokens appear incrementally instead of the UI blocking until the
    full response is ready.  Not suitable for JSON-mode responses (use ``ask``
    with ``json_mode=True`` for those).

    Yields
    ------
    str
        Successive text chunks from the streaming completion.

    Raises
    ------
    RuntimeError
        Same conditions as ``ask()``.
    """
    client = get_client()
    if client is None:
        raise RuntimeError(
            "Groq API key is not set. Paste your key in the sidebar, or add it "
            "to a `.env` file as `GROQ_API_KEY=...` (see README)."
        )

    try:
        stream = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=4096,
            stream=True,
        )
        for chunk in stream:
            delta = chunk.choices[0].delta.content
            if delta:
                yield delta
    except RateLimitError:
        raise RuntimeError(
            "Rate limit reached on your Groq key. Wait a moment and try again."
        )
    except APIConnectionError:
        raise RuntimeError(
            "Could not reach the Groq API. Check your internet connection and try again."
        )
    except APIStatusError as exc:
        raise RuntimeError(f"Groq API error ({exc.status_code}): {exc.message}")


# ---------------------------------------------------------------------------
# File upload -> plain text
# ---------------------------------------------------------------------------

# Accepted file extensions shown in the uploader
SUPPORTED_EXTENSIONS = [
    "pdf", "docx", "doc", "txt",
    "pptx", "ppt",
    "xlsx", "xls", "csv",
    "png", "jpg", "jpeg", "jfif", "webp", "gif", "bmp",
]


def extract_text_from_file(uploaded_file):
    """Extract plain text from an uploaded study document.

    Supported formats
    -----------------
    Text/document : PDF, DOCX/DOC, TXT
    Presentation  : PPTX, PPT
    Spreadsheet   : XLSX, XLS, CSV
    Image         : PNG, JPG, JPEG, JFIF, WEBP, GIF, BMP
        Images are OCR-transcribed via pytesseract (if installed) or
        described via the Groq vision API; falls back to a placeholder if
        neither is available.

    Returns
    -------
    str
        Extracted text, stripped.  Empty string if nothing could be pulled.

    Raises
    ------
    ValueError
        For completely unsupported extensions.
    """
    name = uploaded_file.name.lower()

    # ── PDF ──────────────────────────────────────────────────────────────────
    if name.endswith(".pdf"):
        from pypdf import PdfReader
        reader = PdfReader(uploaded_file)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages).strip()

    # ── DOCX / DOC ───────────────────────────────────────────────────────────
    if name.endswith(".docx") or name.endswith(".doc"):
        import docx
        document = docx.Document(uploaded_file)
        parts = []
        for p in document.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        # Also pull text from tables
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts).strip()

    # ── TXT ──────────────────────────────────────────────────────────────────
    if name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8", errors="ignore").strip()

    # ── PPTX / PPT ───────────────────────────────────────────────────────────
    if name.endswith(".pptx") or name.endswith(".ppt"):
        try:
            from pptx import Presentation
            # Wrap in BytesIO — Streamlit's UploadedFile is not always
            # seekable in the way python-pptx expects a file-path/stream.
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            lines = []
            for slide_num, slide in enumerate(prs.slides, 1):
                lines.append(f"--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
                    # Also pull text from table cells inside shapes
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    lines.append(cell.text.strip())
            return "\n".join(lines).strip()
        except ImportError:
            raise ValueError(
                "python-pptx is not installed. Run: pip install python-pptx"
            )

    # ── XLSX / XLS ───────────────────────────────────────────────────────────
    if name.endswith(".xlsx") or name.endswith(".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(uploaded_file.read()), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join("" if v is None else str(v) for v in row)
                    if row_text.strip():
                        lines.append(row_text)
            return "\n".join(lines).strip()
        except ImportError:
            raise ValueError(
                "openpyxl is not installed. Run: pip install openpyxl"
            )

    # ── CSV ──────────────────────────────────────────────────────────────────
    if name.endswith(".csv"):
        import csv as csv_mod
        text = uploaded_file.read().decode("utf-8", errors="ignore")
        lines = []
        for row in csv_mod.reader(text.splitlines()):
            lines.append("\t".join(row))
        return "\n".join(lines).strip()

    # ── Images (PNG / JPG / JPEG / JFIF / WEBP / GIF / BMP) ─────────────────
    img_exts = (".png", ".jpg", ".jpeg", ".jfif", ".webp", ".gif", ".bmp")
    if any(name.endswith(ext) for ext in img_exts):
        return _extract_text_from_image(uploaded_file)

    raise ValueError(
        f"Unsupported file type '{name.rsplit('.', 1)[-1]}'. "
        f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
    )


# Common Windows install paths for the Tesseract binary.
# pytesseract needs to know where the .exe is when it is not on PATH.
_TESSERACT_CANDIDATE_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    r"C:\Users\Public\tesseract\tesseract.exe",
]


def _configure_tesseract():
    """Point pytesseract at the Tesseract binary if it is not on PATH."""
    try:
        import pytesseract
        # If the binary is already on PATH this will succeed immediately.
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        pass

    # Try well-known Windows install locations.
    import os as _os
    for candidate in _TESSERACT_CANDIDATE_PATHS:
        if _os.path.isfile(candidate):
            try:
                import pytesseract
                pytesseract.pytesseract.tesseract_cmd = candidate
                pytesseract.get_tesseract_version()  # verify it works
                return True
            except Exception:
                pass
    return False


def _extract_text_from_image(uploaded_file):
    """Extract text from an image using Tesseract OCR via pytesseract.

    Automatically locates the Tesseract binary on Windows even when it is
    not on the system PATH.  Returns extracted text, or empty string if
    OCR is unavailable.
    """
    img_bytes = uploaded_file.read()

    try:
        import pytesseract
        from PIL import Image

        _configure_tesseract()

        img = Image.open(io.BytesIO(img_bytes))
        # Convert to RGB so Tesseract handles all formats (JFIF, WEBP, etc.)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img)
        if text.strip():
            return text.strip()
    except Exception:
        pass

    return ""


# ---------------------------------------------------------------------------
# Local ML feature: TF-IDF based key-topic extraction (no API call needed)
# ---------------------------------------------------------------------------

def extract_keywords(text, top_n=10):
    """Extract the most important keywords from *text* using TF-IDF.

    Splits the text into pseudo-sentences (on newlines and full stops) to give
    TF-IDF something to compute inverse-document-frequency over.  Falls back
    to treating the whole text as one document if it is too short to split.

    Parameters
    ----------
    text : str
        The raw notes text.
    top_n : int
        Maximum number of keywords to return.

    Returns
    -------
    list[str]
        Up to *top_n* keywords, ordered by descending TF-IDF score.
        Returns an empty list if the text contains fewer than 2 meaningful
        tokens after stop-word removal (e.g. single-word input).
    """
    sentences = [s.strip() for s in text.replace("\n", ". ").split(".") if s.strip()]
    if len(sentences) < 2:
        sentences = [text]

    vectorizer = TfidfVectorizer(stop_words="english", max_features=200)
    try:
        matrix = vectorizer.fit_transform(sentences)
    except ValueError:
        # Raised when the vocabulary is empty (e.g. text is only stop-words or
        # punctuation) — return an empty list instead of crashing.
        return []

    scores = matrix.sum(axis=0).A1
    terms = vectorizer.get_feature_names_out()
    ranked = sorted(zip(terms, scores), key=lambda x: x[1], reverse=True)
    return [t for t, _ in ranked[:top_n]]
